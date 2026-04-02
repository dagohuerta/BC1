from __future__ import annotations

import json
import tempfile
from pathlib import Path

import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
import streamlit as st

from retail_roi_model.engine import WorkbookDrivenRetailROIModel, outputs_to_jsonable, npv, irr, payback_period
from db_manager import DatabaseManager

import io
import matplotlib.pyplot as plt
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    pass

# --- Inicialización de Base de Datos ---
db = DatabaseManager()
db.init_db()

# --- Carga de Estado de Base de Datos ---
db_modules = db.load_modules()
db_profiles = db.load_profiles()
db_params = db.load_benefit_params()
db_investments = db.load_investments()
db_aspects = db.load_aspect_ranges()

st.set_page_config(page_title="Strategic ROI Analysis", layout="wide")

def format_usd(val):
    if val is None: return "N/D"
    return f"${val:,.0f}" if abs(val) >= 1 else f"${val:,.2f}"

def format_m(val):
    if val is None: return "N/D"
    if abs(val) >= 1_000_000:
        return f"${val/1_000_000:,.1f} M"
    return format_usd(val)

def local_css(file_name):
    with open(file_name) as f:
        st.markdown(f"<style>{f.read()}</style>", unsafe_allow_html=True)

# Cargar Estilos Externos
try:
    local_css("style.css")
except FileNotFoundError:
    pass



def get_total_manual_investment():
    # Sumar inversiones hasta el horizonte seleccionado
    horizon = st.session_state.get('manual_horizon', 5)
    total = 0.0
    for i in range(1, horizon + 1):
        inv = st.session_state.annual_investments.get(i, {'software': 0, 'impl': 0, 'extra': 0})
        total += inv['software'] + inv['impl'] + inv['extra']
    return total

# --- Función para mostrar el Reporte Ejecutivo ---
@st.dialog("Inversiones", width="large")
def configure_investments():
    horizon = st.session_state.get('manual_horizon', 5)
    st.write(f"Ingresa los montos de inversión proyectados para el horizonte de **{horizon} años**.")
    
    # Vista Transpuesta: Años en columnas horizontales
    cols_def = [1.5] + [1] * horizon
    
    # Fila de Encabezados (Años)
    header_cols = st.columns(cols_def)
   # header_cols[0].write("**Categoría de Inversión**") no funciona AV
    for i in range(1, horizon + 1):
        header_cols[i].write(f"**Año {i}**")
        header_cols[i].align = "center"
    #st.divider() - Eliminado por av
        
    # Fila 1: Licencias
    s_cols = st.columns(cols_def)
    s_cols[0].write("Licencias (Software)")
    for i in range(1, horizon + 1):
        val = float(st.session_state.annual_investments.get(i, {}).get('software', 0.0))
        s_cols[i].number_input(f"Soft Y{i}", min_value=0.0, value=val, label_visibility="collapsed", key=f"cfg_soft_{i}")
        
    # Fila 2: Implementación
    i_cols = st.columns(cols_def)
    i_cols[0].write("Implementación")
    for i in range(1, horizon + 1):
        val = float(st.session_state.annual_investments.get(i, {}).get('impl', 0.0))
        i_cols[i].number_input(f"Impl Y{i}", min_value=0.0, value=val, label_visibility="collapsed", key=f"cfg_impl_{i}")
        
    # Fila 3: Adicional / Otros
    e_cols = st.columns(cols_def)
    e_cols[0].write("Adicional / Otros")
    for i in range(1, horizon + 1):
        val = float(st.session_state.annual_investments.get(i, {}).get('extra', 0.0))
        e_cols[i].number_input(f"Extra Y{i}", min_value=0.0, value=val, label_visibility="collapsed", key=f"cfg_extra_{i}")
    
    st.write("")
    if st.button("Guardar en Base de Datos Principal y Cerrar", type="primary", use_container_width=True):
        # Forzamos la captura explícita desde los Session State keys antes de enviar a db_manager
        # Esto soluciona la falla de registro en cache de forms/dialogos
        for i in range(1, horizon + 1):
            if i not in st.session_state.annual_investments:
                st.session_state.annual_investments[i] = {}
            st.session_state.annual_investments[i]['software'] = float(st.session_state.get(f"cfg_soft_{i}", 0.0))
            st.session_state.annual_investments[i]['impl'] = float(st.session_state.get(f"cfg_impl_{i}", 0.0))
            st.session_state.annual_investments[i]['extra'] = float(st.session_state.get(f"cfg_extra_{i}", 0.0))
            
        persist_state() # Guarda toda la maestra configuración
        st.rerun()

@st.dialog("Reporte Ejecutivo de ROI", width="large")
def show_executive_report(results):
    # Definir paleta Dark Mode Neon
    neon_palette = ['#00a3ff', '#00e676', '#ffc107', '#ff5252', '#bb86fc', '#03dac6', '#cf6679']

    # --- Reporte Ejecutivo Estilo Apple (Pixel-Perfect) ---
    st.markdown(f'<div class="header-monday" style="text-align: center; margin-bottom: 5px;">Análisis de Valor Estratégico</div>', unsafe_allow_html=True)
    st.markdown(f'<div style="text-align: center; color: var(--apple-text-dim); margin-bottom: 50px; font-weight: 300; font-size: 1rem; letter-spacing: 0.05em;">{results["cliente"].upper()} &nbsp; • &nbsp; {results["retailer_type"].upper()}</div>', unsafe_allow_html=True)
    
    # Métricas Principales (Apple Product Cards style)
    total_b = sum(results['total_benefit'])
    total_i = sum(results['total_investment'])
    roi_pct = ((total_b - total_i) / total_i * 100) if total_i > 0 else 0
    
    m1, m2, m3, m4, m5 = st.columns(5)
    with m1:
        st.markdown(f"""<div class="monday-card" style="text-align: center;"><div class="monday-card-label">Net Present Value</div><div class="monday-card-value">{format_m(results['npv'])}</div></div>""", unsafe_allow_html=True)
    with m2:
        st.markdown(f"""<div class="monday-card" style="text-align: center;"><div class="monday-card-label">ROI Total</div><div class="monday-card-value" style="color: #00e676;">{roi_pct:,.0f}%</div></div>""", unsafe_allow_html=True)
    with m3:
        st.markdown(f"""<div class="monday-card" style="text-align: center;"><div class="monday-card-label">Internal Rate</div><div class="monday-card-value" style="color: var(--apple-blue);">{results['irr'] if results['irr'] is not None else 'N/D'}%</div></div>""", unsafe_allow_html=True)
    with m4:
        st.markdown(f"""<div class="monday-card" style="text-align: center;"><div class="monday-card-label">Payback Period</div><div class="monday-card-value">{results['payback'] or 'N/D'}</div></div>""", unsafe_allow_html=True)
    with m5:
        st.markdown(f"""<div class="monday-card" style="text-align: center;"><div class="monday-card-label">Total Potential</div><div class="monday-card-value">{format_m(total_b)}</div></div>""", unsafe_allow_html=True)

    st.write("")
    st.write("")

    # Layout de Gráficos (Modern Minimalist)
    g1, g2 = st.columns([1.6, 1])
    
    with g1:
        st.markdown('<div style="font-weight: 500; font-size: 1.1rem; margin-bottom: 25px; color: var(--apple-text-dim);">FLUJO DE CAJA Y RETORNO ESTRATÉGICO</div>', unsafe_allow_html=True)
        years_list = list(range(1, results['years'] + 1))
        
        # Calcular inversiones en negativo para la visualización
        inv_negativas = [-inv for inv in results['total_investment']]
        
        # Calcular flujo de caja descontado para ilustrar trayectoria hacia el NPV
        tasa_desc = float(st.session_state.get('input_discount_rate', 10.0)) / 100
        dcf_acumulado = []
        dcf_actual = 0
        for y, (b, i) in enumerate(zip(results['total_benefit'], results['total_investment']), 1):
            flujo_neto = b - i
            flujo_neto_descontado = flujo_neto / ((1 + tasa_desc) ** y)
            dcf_actual += flujo_neto_descontado
            dcf_acumulado.append(dcf_actual)
        
        fig_traj = go.Figure()
        
        # Barras de Inversión (hacia abajo)
        fig_traj.add_trace(go.Bar(
            x=years_list, 
            y=inv_negativas, 
            name='Inversión Anual', 
            marker_color='#ff5252', 
            opacity=0.85,
            hovertemplate='Año %{x}<br>Inversión: $%{y:,.0f}<extra></extra>'
        ))

        # Barras de Beneficio (hacia arriba)
        fig_traj.add_trace(go.Bar(
            x=years_list, 
            y=results['total_benefit'], 
            name='Beneficio Bruto', 
            marker_color='#00e676',
            opacity=0.85,
            hovertemplate='Año %{x}<br>Beneficio: $%{y:,.0f}<extra></extra>'
        ))

        # Línea de Flujo Acumulado Nominal
        fig_traj.add_trace(go.Scatter(
            x=years_list, 
            y=results['total_cumulative'], 
            name='Flujo Neto Nominal (Sin descontar)', 
            mode='lines+markers+text',
            text=["" for _ in range(len(years_list)-1)] + [f"Nominal: ${results['total_cumulative'][-1]:,.0f}"],
            textposition="top left",
            textfont=dict(color='#2997ff', size=12, family="Arial"),
            line=dict(color='#2997ff', width=2, dash='dot'), 
            marker=dict(size=7, color='white', line=dict(width=1, color='#2997ff')),
            hovertemplate='Año %{x}<br>Acumulado Nominal: $%{y:,.0f}<extra></extra>'
        ))

        # Línea NPV (Flujo Descontado Acumulado) - Clave para un financiero
        fig_traj.add_trace(go.Scatter(
            x=years_list, 
            y=dcf_acumulado, 
            name='NPV Acumulado', 
            mode='lines+markers+text',
            text=["" for _ in range(len(years_list)-1)] + [f"NPV: ${dcf_acumulado[-1]:,.0f}"],
            textposition="top left",
            textfont=dict(color='#ffc107', size=13, family="Arial"),
            line=dict(color='#ffc107', width=4), 
            marker=dict(size=10, color='white', line=dict(width=2, color='#ffc107')),
            hovertemplate='Año %{x}<br>NPV Acumulado: $%{y:,.0f}<extra></extra>'
        ))
        
        # Layout profesional
        fig_traj.update_layout(
            barmode='group',
            template='plotly_dark',
            paper_bgcolor='rgba(0,0,0,0)', 
            plot_bgcolor='rgba(0,0,0,0)', 
            height=420, 
            margin=dict(l=0, r=0, t=0, b=0), 
            legend=dict(orientation="h", y=1.12, x=0, font=dict(size=11)),
            hovermode="x unified"
        )
        
        fig_traj.update_xaxes(showgrid=False, zeroline=False, tickmode='linear', dtick=1)
        # Rango inteligente para el eje Y: Forzar algo de espacio negativo si es muy pequeño para preservar el comparativo
        todos_los_valores = results['total_benefit'] + inv_negativas + results['total_cumulative'] + dcf_acumulado
        max_v = max(todos_los_valores) if todos_los_valores else 1000
        min_v = min(todos_los_valores) if todos_los_valores else -1000
        # Asegurarnos de que haya al menos un 30% de espacio visual abajo si el beneficio es muy grande y la inversión chica.
        if abs(min_v) < (max_v * 0.3):
            rango_min = -(max_v * 0.3)
        else:
            rango_min = min_v * 1.1
            
        fig_traj.update_yaxes(
            showgrid=True, 
            gridcolor='rgba(255,255,255,0.05)', 
            tickprefix="$", 
            zeroline=True,
            zerolinewidth=2,
            zerolinecolor='rgba(255,255,255,0.6)',
            range=[rango_min, max_v * 1.1]
        )
        st.plotly_chart(fig_traj, width="stretch")

    with g2:
        if results['selection']:
            st.markdown('<div style="font-weight: 500; font-size: 1.1rem; margin-bottom: 25px; color: var(--apple-text-dim);">DESGLOSE DE VALOR</div>', unsafe_allow_html=True)
            benefits_by_mod = [sum(results['module_results'][m]['benefit']) for m in results['selection']]
            fig_pie = px.pie(values=benefits_by_mod, names=results['selection'], hole=0.75)
            # Apple Palette: Space Gray shades + Accent Blue
            apple_palette = ['#2997ff', '#323232', '#424245', '#636366', '#86868b']
            fig_pie.update_traces(marker=dict(colors=apple_palette), textinfo='none', hovertemplate='%{label}<br>$%{value:,.0f}')
            fig_pie.update_layout(template='plotly_dark', paper_bgcolor='rgba(0,0,0,0)', showlegend=True, legend=dict(orientation="v", y=0.5, x=1.0), margin=dict(l=0, r=0, t=0, b=0), height=380)
            st.plotly_chart(fig_pie, width="stretch")

    st.write("---")
    st.markdown('<div style="font-weight: 500; font-size: 1.1rem; margin-bottom: 25px; color: var(--apple-text-dim);">DETALLE FINANCIERO ANUAL</div>', unsafe_allow_html=True)
    

    with st.expander("📊 Ver Detalle Anual Proyectado"):
        st.write("Cifras proyectadas año con año (métricas transpuestas)")
        # Lógica de construcción de tabla
        years_list = list(range(1, results['years'] + 1))
        detail_df = pd.DataFrame({
            'Beneficio': results['total_benefit'],
            'Inversión': results['total_investment'],
            'Flujo neto': results.get('total_cashflow', [b - i for b, i in zip(results['total_benefit'], results['total_investment'])]),
            'Cumulado': results['total_cumulative']
        }, index=[f"Año {i}" for i in years_list])
        
        df_t = detail_df.T
        df_t['Total'] = df_t.sum(axis=1)
        st.table(df_t.applymap(format_usd))

    st.write("")

    with st.expander("📚 Racionales de Negocio e Impacto Técnico"):
        benefit_rows = []
        for mod_name in results['selection']:
            mod_data = results['module_results'].get(mod_name, {})
            impacts = []
            if 'aspect_pcts' in results and mod_name in results['aspect_pcts']:
                for aspect, pct in results['aspect_pcts'][mod_name].items():
                    impacts.append(f"{aspect.replace('_',' ').title()}: {pct}%")
            racional = "Calculado según parámetros de industria."
            if mod_name in st.session_state.module_profiles:
                racional = " ".join([st.session_state.aspect_ranges.get(a, {}).get("rationale", "") for a in st.session_state.module_profiles[mod_name]])
            total_monto = sum(mod_data.get('benefit', [0]))
            benefit_rows.append({
                "Módulo": mod_name,
                "Justificación": racional,
                "Impacto": ", ".join(impacts),
                "Total (USD)": format_usd(total_monto)
            })
        st.table(pd.DataFrame(benefit_rows))

    # --- Exportación PPTX Estilo Apple Keynote ---
    try:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        # Colores Apple Keynote
        PP_BG = RGBColor(10, 10, 10)
        PP_ACCENT = RGBColor(0, 113, 227) # Apple Blue
        PP_GRAY = RGBColor(134, 134, 139)
        PP_WHITE = RGBColor(255, 255, 255)
        
        # 1. Slide: Portada Minimalista
        slide_cover = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide_cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PP_BG
        rect.line.fill.background()
        
        title_box = slide_cover.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Strategic ROI Impact"
        p.font.size = Pt(48)
        p.font.name = 'Arial Light' # Intentando estilo fino
        p.font.color.rgb = PP_WHITE
        p.alignment = PP_ALIGN.LEFT
        
        sub_box = slide_cover.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
        p = sub_box.text_frame.paragraphs[0]
        p.text = f"{results['cliente']} | Case Study 2024"
        p.font.size = Pt(18)
        p.font.color.rgb = PP_GRAY
        p.alignment = PP_ALIGN.LEFT

        # 2. Slide: Resumen Ejecutivo (Clean Grid)
        slide_kpi = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide_kpi.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PP_BG
        
        header = slide_kpi.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        p = header.text_frame.paragraphs[0]
        p.text = "Financial Summary"
        p.font.size = Pt(32)
        p.font.color.rgb = PP_WHITE
        
        kpi_list = [
            ("NET PRESENT VALUE", format_m(results['npv'])),
            ("INTERNAL RATE", f"{results['irr']}%"),
            ("PAYBACK PERIOD", results['payback']),
            ("TOTAL POTENTIAL", format_m(sum(results['total_benefit'])))
        ]
        
        left, top = Inches(1), Inches(2)
        for name, val in kpi_list:
            # Label
            lbl = slide_kpi.shapes.add_textbox(left, top, Inches(4), Inches(0.3))
            p = lbl.text_frame.paragraphs[0]
            p.text = name
            p.font.size = Pt(10)
            p.font.color.rgb = PP_GRAY
            
            # Value
            vbox = slide_kpi.shapes.add_textbox(left, top + Inches(0.3), Inches(4), Inches(0.7))
            p = vbox.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(36)
            p.font.color.rgb = PP_WHITE
            
            top += Inches(1.2)
            if top > Inches(5):
                top = Inches(2)
                left = Inches(5.5)

        # 3. Slide: Visual Data (Minimalist Chart)
        slide_chart = prs.slides.add_slide(prs.slide_layouts[6])
        rect = slide_chart.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = PP_BG
        
        plt.style.use('dark_background')
        fig, ax = plt.subplots(figsize=(10, 5), facecolor='#0a0a0a')
        ax.set_facecolor('#0a0a0a')
        years = list(range(1, results['years'] + 1))
        ax.bar(years, results['total_benefit'], color='white', alpha=0.1, width=0.6)
        ax.plot(years, results['total_cumulative'], color='#0071e3', linewidth=3, marker='o', markersize=8)
        ax.set_title("Cash Flow Path", color='#86868b', fontsize=14, loc='left')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', color='#1c1c1e', alpha=0.5)
        
        img_buf = io.BytesIO()
        fig.savefig(img_buf, format='png', bbox_inches='tight', dpi=150, facecolor=fig.get_facecolor())
        img_buf.seek(0)
        plt.close(fig)
        slide_chart.shapes.add_picture(img_buf, Inches(0.5), Inches(1.5), width=Inches(9))

        # Guardar en memoria
        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        pptx_buf.seek(0)

        st.write("---")
        c1, c2 = st.columns([3, 1])
        with c2:
            st.download_button(
                "📈 EXPORT KEYNOTE STYLE",
                data=pptx_buf,
                file_name=f"ROI_Presentation_{results['cliente']}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                help="Exportar a PowerPoint con estética minimalista premium"
            )
    except Exception as e:
        st.warning(f"Exportación no disponible: {e}")

    st.info("Este reporte resume los beneficios proyectados basados en los parámetros ingresados y las curvas de adopción estándar de la industria.")

# Rangos por aspecto - inicializar en session_state
if 'aspect_ranges' not in st.session_state:
    if db_aspects:
        st.session_state.aspect_ranges = db_aspects
    else:
        st.session_state.aspect_ranges = {
            "sales": {"min": 0, "max": 20, "rationale": "Optimización de ingresos mediante mejor disponibilidad y previsión operativa coordinada."},
            "inventory_reduction": {"min": 0, "max": 15, "rationale": "Reducción de capital inmovilizado y mejora en eficiencia de stock mediante visibilidad."},
            "margin": {"min": 0, "max": 10, "rationale": "Protección y expansión de margen mediante optimización estratégica de precios y promociones."},
            "labor": {"min": 0, "max": 12, "rationale": "Aumento de productividad del personal y eficiencia en operaciones de tienda/almacén."},
            "logistics": {"min": 0, "max": 8, "rationale": "Optimización de costos logísticos y de transporte mediante una red inteligente."},
        }

# Inicializar inversiones anuales en session_state si no existe
if 'module_options' not in st.session_state:
    # Cargar de DB si hay datos, si no usar defaults
    if db_modules:
        st.session_state.module_options = db_modules
    else:
        st.session_state.module_options = [
            "Inventory Optimization", "Pricing", "Merchandising", "Customer Experience",
            "Supply Chain", "Retail Insights", "Store Operations", "Loyalty", "Omnichannel", "Data Science"
        ]

if 'module_profiles' not in st.session_state:
    if db_profiles:
        st.session_state.module_profiles = db_profiles
    else:
        st.session_state.module_profiles = {
            "Inventory Optimization": ["inventory_reduction"],
            "Pricing": ["margin"],
            "Merchandising": ["sales"],
            "Customer Experience": ["sales", "margin"],
            "Supply Chain": ["logistics"],
            "Retail Insights": ["sales", "inventory_reduction"],
            "Store Operations": ["labor"],
            "Loyalty": ["sales"],
            "Omnichannel": ["sales", "logistics"],
            "Data Science": ["margin", "inventory_reduction"],
        }

if 'benefit_params' not in st.session_state:
    if db_params:
        st.session_state.benefit_params = db_params
    else:
        st.session_state.benefit_params = {}
        for module in st.session_state.module_options:
            st.session_state.benefit_params[module] = {}
            for aspect in st.session_state.module_profiles.get(module, []):
                cfg = st.session_state.aspect_ranges[aspect]
                min_val, max_val = cfg["min"], cfg["max"]
                st.session_state.benefit_params[module][aspect] = {"min": min_val + 1, "max": max_val - 1}

if 'annual_investments' not in st.session_state:
    if db_investments:
        st.session_state.annual_investments = db_investments
    else:
        st.session_state.annual_investments = {i: {'software': 0, 'impl': 0, 'extra': 0} for i in range(1, 11)}

# Helper para guardar estado
def persist_state():
    db.sync_all(
        st.session_state.module_options,
        st.session_state.module_profiles,
        st.session_state.benefit_params,
        st.session_state.annual_investments,
        st.session_state.aspect_ranges
    )

aspect_ranges = st.session_state.aspect_ranges

# --- Estado de la Aplicación (Dashboard vs Editor) ---
if 'app_mode' not in st.session_state:
    st.session_state.app_mode = "dashboard"

if 'editing_exercise_id' not in st.session_state:
    st.session_state.editing_exercise_id = None

if 'editing_exercise_name' not in st.session_state:
    st.session_state.editing_exercise_name = ""

@st.dialog("Guardar Ejercicio de ROI", width="small")
def save_exercise_dialog():
    is_update = st.session_state.editing_exercise_id is not None
    st.write("Introduce un nombre para identificar este análisis." if not is_update else "Actualiza el nombre si lo deseas.")
    
    default_name = st.session_state.editing_exercise_name if is_update else ""
    e_name = st.text_input("Nombre del Ejercicio", value=default_name, placeholder="Ej: Caso Walmart 2024 - Escenario Optimista")
    
    btn_label = "Confirmar Actualización" if is_update else "Confirmar Guardado"
    
    if st.button(btn_label, type="primary", use_container_width=True):
        if e_name:
            # Recopilar todos los inputs actuales desde las claves de widgets
            data_to_save = {
                'exercise_name': e_name,
                'client_name': st.session_state.get('input_client_name', "Cliente Demo"),
                'retailer_type': st.session_state.get('input_retailer_type', "Retail general"),
                'net_revenue': st.session_state.get('input_net_revenue', 100000000.0),
                'growth_rate': st.session_state.get('input_growth_rate', 3.0),
                'inventory': st.session_state.get('input_inventory', 25000000.0),
                'carrying_cost': st.session_state.get('input_carrying_cost', 20.0),
                'cogs_pct': st.session_state.get('input_cogs_pct', 55.0),
                'sga_pct': st.session_state.get('input_sga_pct', 15.0),
                'tax_rate': st.session_state.get('input_tax_rate', 30.0),
                'discount_rate': st.session_state.get('input_discount_rate', 10.0),
                'adoption_years': st.session_state.get('manual_horizon', 5),
                'scenario_type': st.session_state.get('input_scenario_mode', "Base"),
                'module_selected': st.session_state.get('input_module_selected', []),
                'module_benefits': st.session_state.get('module_benefits', {}), 
                'annual_investments': st.session_state.annual_investments
            }
            
            if is_update:
                data_to_save['id'] = st.session_state.editing_exercise_id
            
            saved_id = db.save_exercise(data_to_save)
            if not is_update:
                st.session_state.editing_exercise_id = saved_id
            st.session_state.editing_exercise_name = e_name
            
            st.success(f"Ejercicio '{e_name}' {'actualizado' if is_update else 'guardado'} correctamente.")
            st.rerun()
        else:
            st.error("Por favor, ingresa un nombre.")

def load_exercise_into_state(e_id):
    ex = db.load_exercise(e_id)
    if ex:
        # Forzar el modo de entrada a manual para que se vean los widgets
        st.session_state.input_entry_mode = "Entrada manual"
        
        # Actualizar claves de widgets directamente para forzar renderizado
        st.session_state.input_client_name = str(ex['client_name'])
        st.session_state.input_retailer_type = str(ex['retailer_type'])
        st.session_state.input_net_revenue = float(ex['net_revenue'])
        st.session_state.input_growth_rate = float(ex['growth_rate'])
        st.session_state.input_inventory = float(ex['inventory'])
        st.session_state.input_carrying_cost = float(ex['carrying_cost'])
        st.session_state.input_cogs_pct = float(ex['cogs_pct'])
        st.session_state.input_sga_pct = float(ex['sga_pct'])
        st.session_state.input_tax_rate = float(ex['tax_rate'])
        st.session_state.input_discount_rate = float(ex['discount_rate'])
        st.session_state.manual_horizon = int(ex['adoption_years'])
        st.session_state.input_scenario_mode = str(ex['scenario_type'])
        st.session_state.input_module_selected = ex['module_selected']
        
        # También actualizar variables de cálculo
        st.session_state.annual_investments = ex['annual_investments']
        st.session_state.editing_exercise_id = e_id
        st.session_state.editing_exercise_name = ex['exercise_name']
        st.session_state.app_mode = "editor"
        st.rerun()

# --- VISTA: DASHBOARD ---
if st.session_state.app_mode == "dashboard":
    st.markdown('<div class="header-monday" style="text-align:center; font-size: 2.5rem; margin-top: 50px;">Bienvenido al ROI Strategist</div>', unsafe_allow_html=True)
    st.markdown('<div style="text-align:center; color: var(--apple-text-dim); margin-bottom: 50px; font-weight: 300;">Gestión Ejecutiva de Proyecciones de Retorno de Inversión</div>', unsafe_allow_html=True)
    
    col_dash1, col_dash2 = st.columns(2)
    
    with col_dash1:
        with st.container(border=True):
            st.markdown('<div style="font-size: 1.5rem; font-weight: 500; margin-bottom: 10px;">🆕 Nuevo Análisis</div>', unsafe_allow_html=True)
            st.write("Comienza un nuevo ejercicio de ROI desde cero utilizando los parámetros maestros.")
            if st.button("Crear Nuevo Ejercicio", type="primary", use_container_width=True):
                # Resetear variables de entrada para un nuevo ejercicio
                st.session_state.input_entry_mode = "Entrada manual"
                st.session_state.input_client_name = "Nuevo Cliente"
                st.session_state.input_net_revenue = 100000000.0
                st.session_state.input_growth_rate = 3.0
                st.session_state.input_inventory = 25000000.0
                st.session_state.input_module_selected = st.session_state.module_options[:2]
                st.session_state.editing_exercise_id = None
                st.session_state.editing_exercise_name = ""
                st.session_state.app_mode = "editor"
                st.rerun()
                
    with col_dash2:
        with st.container(border=True):
            st.markdown('<div style="font-size: 1.5rem; font-weight: 500; margin-bottom: 10px;">🗂️ Cargar Existente</div>', unsafe_allow_html=True)
            st.write("Abre un ejercicio guardado previamente para revisar o editar.")
            saved_exercises = db.get_exercise_list()
            if not saved_exercises:
                st.info("No tienes ejercicios guardados aún.")
            else:
                for eid, ename, eclient, edate in saved_exercises:
                    c_load1, c_load2 = st.columns([4, 1])
                    with c_load1:
                        st.markdown(f"**{ename}**")
                        st.caption(f"{eclient} • {edate}")
                    with c_load2:
                        if st.button("Abrir", key=f"load_{eid}", use_container_width=True):
                            load_exercise_into_state(eid)
                    st.divider()

    st.stop() # Detener ejecución aquí si estamos en dashboard

# --- VISTA: EDITOR ---
st.title("Strategic ROI Analyst")

main_tab, config_tab = st.tabs(["Cálculo", "Configuración"])

with config_tab:
    config_subtab1, config_subtab2, config_subtab3, config_subtab4 = st.tabs(["Módulos Oracle", "Parámetros de Beneficio", "Aspectos", "Fórmulas Financieras"])

    with config_subtab1:
        st.subheader("Gestión de Módulos Oracle")
        
        # --- Toolbar Dual-View (Agregar + Buscar + Vista) ---
        with st.container():
            t1, t2, t3 = st.columns([2, 2, 1])
            with t1:
                new_module = st.text_input("Nombre del módulo", key="new_module_compact", label_visibility="collapsed", placeholder="➕ Nuevo módulo Oracle...")
                if new_module and new_module not in st.session_state.module_options:
                    if st.button("Crear Módulo", key="btn_create_dual"):
                        st.session_state.module_options.append(new_module)
                        st.session_state.module_profiles[new_module] = []
                        st.session_state.benefit_params[new_module] = {}
                        persist_state()
                        st.rerun()
            with t2:
                search_query = st.text_input("🔍 Buscar...", key="module_search_dual", label_visibility="collapsed", placeholder="🔍 Buscar módulos...")
            with t3:
                view_mode = st.radio("Vista", ["⊞", "☰"], index=1, horizontal=True, label_visibility="collapsed")

        # Mapeo de modo de vista
        grid_mode = (view_mode == "⊞")
        filtered_modules = [m for m in st.session_state.module_options if search_query.lower() in m.lower()]
        
        if not filtered_modules:
            st.warning("No se encontraron resultados.")
        else:
            if grid_mode:
                # Grid Optimizado (3 columnas)
                cols_per_row = 3
                for i in range(0, len(filtered_modules), cols_per_row):
                    row_modules = filtered_modules[i:i + cols_per_row]
                    cols = st.columns(cols_per_row)
                    for j, module in enumerate(row_modules):
                        with cols[j]:
                            with st.container(border=True):
                                n_aspects = len(st.session_state.module_profiles.get(module, []))
                                st.markdown(f"**{module}** <small style='color:var(--dark-text-dim); float:right;'>{n_aspects} asp.</small>", unsafe_allow_html=True)
                                
                                current_aspects = st.session_state.module_profiles.get(module, [])
                                selected_aspects = st.multiselect("Impacto", list(aspect_ranges.keys()), default=current_aspects, key=f"aspects_grid_{module}", label_visibility="collapsed")
                                
                                if selected_aspects != current_aspects:
                                    st.session_state.module_profiles[module] = selected_aspects
                                    for aspect in selected_aspects:
                                        if aspect not in st.session_state.benefit_params[module]:
                                            min_v, max_v = aspect_ranges[aspect]
                                            st.session_state.benefit_params[module][aspect] = {"min": min_v + 1, "max": max_v - 1}
                                    for a in list(st.session_state.benefit_params[module].keys()):
                                        if a not in selected_aspects:
                                            del st.session_state.benefit_params[module][a]
                                    persist_state()
                                    st.rerun()

                                with st.expander("Opciones", expanded=False):
                                    new_name = st.text_input("Nombre", module, key=f"rename_grid_{module}", label_visibility="collapsed")
                                    if new_name != module and new_name.strip() != "":
                                        idx = st.session_state.module_options.index(module)
                                        st.session_state.module_options[idx] = new_name
                                        st.session_state.module_profiles[new_name] = st.session_state.module_profiles.pop(module)
                                        st.session_state.benefit_params[new_name] = st.session_state.benefit_params.pop(module)
                                        st.rerun()
                                    if st.button("🗑️ Eliminar", key=f"del_grid_{module}", use_container_width=True):
                                        st.session_state.module_options.remove(module)
                                        del st.session_state.module_profiles[module]
                                        del st.session_state.benefit_params[module]
                                        persist_state()
                                        st.rerun()
            else:
                # Vista de Lista (Filas compactas)
                for module in filtered_modules:
                    with st.container(border=True):
                        c1, c2, c3 = st.columns([2, 3, 1])
                        with c1:
                            st.markdown(f"**{module}**")
                            n_aspects = len(st.session_state.module_profiles.get(module, []))
                            st.caption(f"{n_aspects} aspectos seleccionados")
                        with c2:
                            current_aspects = st.session_state.module_profiles.get(module, [])
                            selected_aspects = st.multiselect("Impacto", list(aspect_ranges.keys()), default=current_aspects, key=f"aspects_list_{module}", label_visibility="collapsed")
                            if selected_aspects != current_aspects:
                                st.session_state.module_profiles[module] = selected_aspects
                                for aspect in selected_aspects:
                                    if aspect not in st.session_state.benefit_params[module]:
                                        cfg = aspect_ranges[aspect]
                                        min_v, max_v = cfg["min"], cfg["max"]
                                        st.session_state.benefit_params[module][aspect] = {"min": min_v + 1, "max": max_v - 1}
                                for a in list(st.session_state.benefit_params[module].keys()):
                                    if a not in selected_aspects:
                                        del st.session_state.benefit_params[module][a]
                                st.rerun()
                        with c3:
                            with st.popover("⚙️"):
                                new_name = st.text_input("Renombrar módulo", module, key=f"rename_list_{module}")
                                if new_name != module and new_name.strip() != "":
                                    idx = st.session_state.module_options.index(module)
                                    st.session_state.module_options[idx] = new_name
                                    st.session_state.module_profiles[new_name] = st.session_state.module_profiles.pop(module)
                                    st.session_state.benefit_params[new_name] = st.session_state.benefit_params.pop(module)
                                    st.rerun()
                                if st.button("🗑️ Eliminar Módulo", key=f"del_list_{module}", type="secondary", use_container_width=True):
                                    st.session_state.module_options.remove(module)
                                    del st.session_state.module_profiles[module]
                                    del st.session_state.benefit_params[module]
                                    persist_state()
                                    st.rerun()

    with config_subtab2:
        st.subheader("Parámetros de Beneficio por Módulo y Aspecto")
        for module in st.session_state.module_options:
            st.write(f"**{module}**")
            if module in st.session_state.module_profiles:
                for aspect in st.session_state.module_profiles[module]:
                    col_min, col_max = st.columns(2)
                    cfg = aspect_ranges[aspect]
                    min_val, max_val = cfg["min"], cfg["max"]
                    with col_min:
                        current_min = st.session_state.benefit_params[module][aspect]["min"]
                        new_min = st.number_input(f"Min % {module} - {aspect}", float(min_val), float(max_val), float(current_min), key=f"min_{module}_{aspect}")
                        st.session_state.benefit_params[module][aspect]["min"] = new_min
                    with col_max:
                        current_max = st.session_state.benefit_params[module][aspect]["max"]
                        new_max = st.number_input(f"Max % {module} - {aspect}", float(min_val), float(max_val), float(current_max), key=f"max_{module}_{aspect}")
                        st.session_state.benefit_params[module][aspect]["max"] = new_max
            st.write("---")

    with config_subtab3:
        st.subheader("Configuración de Aspectos e Impacto")
        st.info("Administra el catálogo maestro de aspectos de valor. Configura sus rangos esperados de impacto y su justificación financiera oficial.")
        
        # Agregar nuevo aspecto (Oculto en expander por defecto para limpieza visual)
        with st.expander("➕ Formulario: Crear Nuevo Aspecto de Valor", expanded=False):
            st.write("Registra un nuevo aspecto para asociarlo dinámicamente a tus módulos de negocio en la aplicación.")
            col_new_name, col_new_min, col_new_max = st.columns([2, 1, 1])
            with col_new_name:
                new_aspect_name = st.text_input("Identificador Técnico (e.g. 'marketing_roi')", key="new_aspect_name")
            with col_new_min:
                new_aspect_min = st.number_input("Escenario Conservador (%)", min_value=0.0, max_value=100.0, value=0.0, key="new_aspect_min")
            with col_new_max:
                new_aspect_max = st.number_input("Escenario Optimista (%)", min_value=0.0, max_value=100.0, value=20.0, key="new_aspect_max")
            
            new_aspect_rat = st.text_area("Justificación Comercial / Racional Financiero", key="new_aspect_rat", placeholder="Describe argumentativamente por qué y cómo este aspecto genera un retorno de inversión. Esta descripción alimentará los reportes ejecutivos generados.", height=80)
            
            # Contenedor para alinear el botón
            col_esp, col_add_btn = st.columns([3, 1])
            with col_add_btn:
                if st.button("Guardar en Catálogo", key="btn_add_aspect", type="primary", use_container_width=True):
                    if new_aspect_name and new_aspect_name not in st.session_state.aspect_ranges:
                        if new_aspect_max > new_aspect_min:
                            st.session_state.aspect_ranges[new_aspect_name] = {"min": new_aspect_min, "max": new_aspect_max, "rationale": new_aspect_rat}
                            st.success(f"Aspecto '{new_aspect_name}' integrado.")
                            persist_state()
                            st.rerun()
                        else:
                            st.error("El esquema optimista debe ser > al conservador.")
                    else:
                        st.error("Nombre inválido o el aspecto ya existe.")
        
        st.write("")
        # Listar y editar aspectos
        st.markdown("### 📋 Catálogo Activo de Indicadores de Retorno")
        
        aspect_items = list(st.session_state.aspect_ranges.items())
        cols_per_row = 2
        for i in range(0, len(aspect_items), cols_per_row):
            row_aspects = aspect_items[i:i + cols_per_row]
            cols = st.columns(cols_per_row)
            
            for j, (aspect_name, val_dict) in enumerate(row_aspects):
                with cols[j]:
                    min_val = float(val_dict["min"])
                    max_val = float(val_dict["max"])
                    rat_val = val_dict.get("rationale", "")
                    
                    with st.container(border=True):
                        # Cabecera de Tarjeta
                        c_title, c_del = st.columns([5, 1.5])
                        with c_title:
                            st.markdown(f"<h5 style='margin-top: 5px; margin-bottom: 10px; color: var(--apple-blue); letter-spacing: 0.5px;'>🔹 {str(aspect_name).replace('_', ' ').title()}</h5>", unsafe_allow_html=True)
                        with c_del:
                            if st.button("🗑️ Borrar", key=f"del_aspect_{aspect_name}", help="Eliminar este aspecto para siempre", use_container_width=True):
                                modules_using = [m for m, aspects in st.session_state.module_profiles.items() if aspect_name in aspects]
                                if modules_using:
                                    st.error("En uso.")
                                else:
                                    del st.session_state.aspect_ranges[aspect_name]
                                    st.rerun()
                        
                        # Cuerpo de Tarjeta: Inputs 50/50 y luego el Texto abajo
                        c_pmin, c_pmax = st.columns(2)
                        with c_pmin:
                            new_min = st.number_input(f"Piso (%)", min_value=0.0, max_value=100.0, value=min_val, key=f"emin_{aspect_name}")
                        with c_pmax:
                            new_max = st.number_input(f"Techo (%)", min_value=0.0, max_value=100.0, value=max_val, key=f"emax_{aspect_name}")
                        
                        st.markdown("<div style='font-size:0.8rem; margin-bottom: 5px; color: var(--apple-text-dim);'>Justificación Funcional</div>", unsafe_allow_html=True)
                        new_rat = st.text_area("Racional", value=rat_val, key=f"erat_{aspect_name}", height=68, label_visibility="collapsed", placeholder="Describe la métrica...")
                        
                    # Sincronización Automática
                    if (new_min, new_max, new_rat) != (min_val, max_val, rat_val):
                        if new_max > new_min:
                            st.session_state.aspect_ranges[aspect_name] = {"min": new_min, "max": new_max, "rationale": new_rat}
                            for module, aspects in st.session_state.module_profiles.items():
                                if aspect_name in aspects:
                                    st.session_state.benefit_params[module][aspect_name]["min"] = min(st.session_state.benefit_params[module][aspect_name]["min"], new_max - 0.1) # Safe lock
                                    st.session_state.benefit_params[module][aspect_name]["max"] = min(st.session_state.benefit_params[module][aspect_name]["max"], new_max - 0.1)
                        else:
                            st.error(f"Techo > Piso en {aspect_name}")
            
        st.write("")
        
        if st.button("💾 Guardar Configuración en Base de Datos", type="primary", use_container_width=True):
            persist_state()
            st.success("Configuración sincronizada con la base de datos.")

    with config_subtab4:
        st.subheader("Glosario de Indicadores y Fórmulas")
        st.info("Listado oficial de las métricas clave calculadas por el motor financiero del ROI Strategist, con su metodología y propósito.")
        
        indicators = [
            {
                "Indicador": "Net Present Value (NPV)", 
                "Fórmula Algorítmica": "Σ [ Flujo Neto(t) / (1 + r)^t ] - Inversión Inicial", 
                "Descripción y Propósito": "Valor actual de los flujos de efectivo futuros descontados por la tasa de riesgo (r = tasa de descuento). Indica el valor total de la riqueza creada a día de hoy."
            },
            {
                "Indicador": "ROI Total", 
                "Fórmula Algorítmica": "[(Beneficio Acumulado - Inversión Acumulada) / Inversión Acumulada] * 100", 
                "Descripción y Propósito": "Retorno de Inversión sobre el horizonte del proyecto. Mide la rentabilidad global bruta generada por el capital invertido."
            },
            {
                "Indicador": "Internal Rate of Return (IRR)", 
                "Fórmula Algorítmica": "La tasa 'r' que hace que NPV = 0", 
                "Descripción y Propósito": "Tasa de interés implícita que iguala el valor presente de los beneficios futuros con el costo de inversión. Sirve para comparar contra otras inversiones posibles."
            },
            {
                "Indicador": "Payback Period", 
                "Fórmula Algorítmica": "Año(n) de cruce + [ (Saldo a recuperar) / Flujo Neto (Año n+1) ]", 
                "Descripción y Propósito": "Tiempo estandarizado (expresado en Años y Meses) requerido para recuperar íntegramente la inversión inicial con los beneficios netos operativos."
            },
            {
                "Indicador": "Total Potential (Beneficio Neto)", 
                "Fórmula Algorítmica": "Σ Beneficios Anuales de módulos parametrizados", 
                "Descripción y Propósito": "Representa el impacto comercial bruto total de todos los módulos operando conjuntamente antes del descuento de la inversión."
            },
            {
                "Indicador": "Flujo Neto Creciente Acumulado", 
                "Fórmula Algorítmica": "Σ [ Beneficio(t) - Inversión(t) ] desde el Año 1 hasta t", 
                "Descripción y Propósito": "Balance entre la salida y entrada de efectivo, utilizado para determinar el marco de autofinaciamiento y trazar la gráfica de Break-even."
            }
        ]
        
        st.table(pd.DataFrame(indicators))

with main_tab:
    with st.sidebar:
        st.markdown('<div class="header-monday" style="font-size: 1.6rem; border-bottom: 5px solid var(--apple-blue);">ROI Strategist</div>', unsafe_allow_html=True)
        
        st.write("---")
        if st.button("💾 Guardar Ejercicio", use_container_width=True, type="primary"):
            save_exercise_dialog()
            
        if st.button("🏠 Ir al Dashboard", use_container_width=True):
            st.session_state.app_mode = "dashboard"
            st.rerun()

        st.write("---")
        st.write("**Parámetros del Modelo**")
        st.caption("v5.0.0 - Obsidian Dark UI")
        
        st.write("---")
        st.write("**Configuración de Inversión**")
        if st.button("🛠️ Inversiones Anuales", use_container_width=True):
            configure_investments()
        total_inv_setup = get_total_manual_investment()
        st.caption(f"Total Configurado: {format_usd(total_inv_setup)}")

    st.markdown('<div class="header-monday" style="text-align:center;">Strategic ROI Analyst</div>', unsafe_allow_html=True)
    st.info("Bienvenido. Ejecuta el análisis de retorno de inversión en el nuevo entorno Obsidian Dark.")
    
    mode = st.radio("Método de entrada", ["Carga Excel", "Entrada manual"], horizontal=True, key="input_entry_mode")

    if mode == "Carga Excel":
        st.caption("Carga el archivo .xlsm y genera el resumen de ROI y el JSON consolidado.")
        uploaded = st.file_uploader("Workbook Excel (.xlsm)", type=["xlsm", "xlsx"])

        if uploaded is not None:
            with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{uploaded.name}") as tmp:
                tmp.write(uploaded.getbuffer())
                temp_path = Path(tmp.name)

            if st.button("Calcular ROI"):
                try:
                    model = WorkbookDrivenRetailROIModel(temp_path)
                    outputs = model.run()
                    payload = outputs_to_jsonable(outputs)
                    
                    # Preparar data del reporte
                    report_data = {
                        "cliente": payload["summary_metrics"].get("client_name", "Cliente"),
                        "retailer_type": payload["summary_metrics"].get("retailer_type", "Retailer"),
                        "npv": payload["summary_metrics"].get("npv", 0),
                        "irr": payload["summary_metrics"].get("irr_pct", 0),
                        "payback": payload["summary_metrics"].get("payback_period"),
                        "years": 5,
                        "total_benefit": payload["annual_total"]["total_estimated_benefit"],
                        "total_investment": payload["annual_total"]["total_investment"],
                        "total_cumulative": payload["annual_total"]["cumulative_net_benefit"],
                        "selection": [m["module_name"] for m in payload.get("module_results", []) if m.get("selected")],
                        "module_results": {m["module_name"]: m for m in payload.get("module_results", [])}
                    }
                    # Mostrar reporte inmediatamente
                    st.session_state.excel_results = report_data
                    show_executive_report(report_data)

                except Exception as exc:
                    st.error(f"No fue posible calcular el modelo: {exc}")

        if "excel_results" in st.session_state:
            st.write("---")
            if st.button("📊 Ver Reporte Ejecutivo (Excel)", use_container_width=True):
                show_executive_report(st.session_state.excel_results)
            st.info("Cálculo completado. El detalle completo está disponible en el Reporte Ejecutivo pop-up.")

    if mode == "Entrada manual":
        st.caption("Completa los datos clave y calcula ROI automáticamente sin usar Excel.")
        
        st.subheader("1. Horizonte de análisis")
        st.info("Define el tiempo de retorno. Cambiar este valor actualizará automáticamente el configurador de inversiones.")
        adoption_years = st.radio("Horizonte de análisis (años)", [3, 5], index=1, key="manual_horizon", horizontal=True)

        st.subheader("1. Datos base")
        
        client_name = st.text_input("Cliente", key="input_client_name")
        retailer_types = ["Retail general", "Supermercado", "Moda", "Electrónica"]
        retailer_type = st.selectbox("Tipo de retailer", retailer_types, key="input_retailer_type")

        col1, col2, col3 = st.columns(3)
        with col1:
            net_revenue = st.number_input("Ingresos brutos anuales (USD)", min_value=0.0, step=1_000.0, key="input_net_revenue")
            cogs_pct = st.number_input("% COGS sobre ingresos", min_value=0.0, max_value=100.0, step=0.1, key="input_cogs_pct")
            sga_pct = st.number_input("% SGA sobre ingresos", min_value=0.0, max_value=100.0, step=0.1, key="input_sga_pct")
        with col2:
            tax_rate = st.number_input("% tasa de impuestos", min_value=0.0, max_value=100.0, step=0.1, key="input_tax_rate")
            inventory = st.number_input("Inventario total (USD)", min_value=0.0, step=1000.0, key="input_inventory")
            carrying_cost = st.number_input("% costo de mantenimiento de inventario", min_value=0.0, max_value=100.0, step=0.1, key="input_carrying_cost")
        with col3:
            growth_rate = st.number_input("% crecimiento de ingresos anual", min_value=-100.0, max_value=100.0, step=0.1, key="input_growth_rate")
            discount_rate = st.number_input("% tasa de descuento NPV", min_value=0.0, max_value=100.0, step=0.1, key="input_discount_rate")

        st.subheader("2. Selección de módulos Oracle")
        module_options = st.session_state.module_options
        module_selected = st.multiselect("Selecciona módulos Oracle", module_options, key="input_module_selected")

        module_profiles = st.session_state.module_profiles
        benefit_params = st.session_state.benefit_params

        module_benefits = {}
        if module_selected:
            st.write("Usando parámetros configurados.")
            scenarios = ["Conservative", "Base", "Aggressive"]
            scenario = st.radio("Escenario", scenarios, key="input_scenario_mode")
            for module in module_selected:
                module_benefits[module] = {}
                for aspect in module_profiles.get(module, []):
                    params = benefit_params.get(module, {}).get(aspect, {"min": 5.0, "max": 15.0})
                    if scenario == "Conservative":
                        pct = params["min"]
                    elif scenario == "Base":
                        pct = (params["min"] + params["max"]) / 2
                    else:  # Aggressive
                        pct = params["max"]
                    module_benefits[module][aspect] = pct
        else:
            st.info("Selecciona al menos un módulo para obtener resultados por módulo.")

        st.subheader("3. Inversiones por año")
        st.info("Utiliza el botón '🛠️ Inversiones Anuales' en la barra lateral para ajustar el cronograma multi-anual.")
        
        # Resumen visual discreto
        total_inv_setup = get_total_manual_investment()
        st.caption(f"**Total de Inversión**: {format_usd(total_inv_setup)}")

        # Botón de cálculo
        if st.button("Calcular ROI manual", type="primary", use_container_width=True):
            def calc_manual_roi():
                # o forzar lectura de session_state si hay duda
                l_rev = net_revenue
                l_cogs_p = cogs_pct
                l_sga_p = sga_pct
                l_inv = inventory
                l_carry = carrying_cost
                l_growth = growth_rate
                l_discount = discount_rate
                l_years = adoption_years
                l_selected = module_selected
                
                cogs = l_rev * l_cogs_p / 100.0
                sga = l_rev * l_sga_p / 100.0
                operating_income = l_rev - cogs - sga
                base_margin = (operating_income / l_rev * 100.0) if l_rev > 0 else 0.0

                years = adoption_years
                yearly_total_benefit = []
                yearly_total_investment = []
                yearly_total_cashflow = []
                yearly_total_cumulative = []
                module_results = {m: {"benefit": [], "investment": [], "cashflow": [], "cumulative": []} for m in module_selected}

                # --- CALCULO DINAMICO DE BENEFICIOS ---
                # Identificar qué base financiera afecta cada aspecto (Ventas -> Ingresos, etc.)
                all_aspects = set()
                for mod in module_selected:
                    all_aspects.update(module_benefits.get(mod, {}).keys())
                
                aspect_totals = {a: 0.0 for a in all_aspects}
                for mod in module_selected:
                    for a, pct in module_benefits.get(mod, {}).items():
                        aspect_totals[a] += pct

                cumulative_total = 0.0
                for year in range(1, years + 1):
                    # Financieros base del año
                    rev_year = net_revenue * ((1 + growth_rate / 100.0) ** (year - 1))
                    cogs_year = rev_year * cogs_pct / 100.0
                    sga_year = rev_year * sga_pct / 100.0
                    base_oi_year = rev_year - cogs_year - sga_year

                    inv_data = st.session_state.annual_investments.get(year, {'software': 0, 'impl': 0, 'extra': 0})
                    total_investment = inv_data['software'] + inv_data['impl'] + inv_data['extra']

                    # Calcular beneficios totales del año recorriendo todos los aspectos dinámicamente
                    total_benefit = 0.0
                    for aspect, total_pct in aspect_totals.items():
                        a_lower = aspect.lower()
                        # Determinar base de impacto
                        if any(k in a_lower for k in ["inv", "stock"]):
                            base = inventory * (carrying_cost / 100.0)
                        elif any(k in a_lower for k in ["lab", "sga", "pers", "tra"]):
                            base = sga_year
                        elif any(k in a_lower for k in ["marg", "prof", "log"]):
                            base = base_oi_year
                        else:
                            base = rev_year # Default: Sales/Revenue impact
                        
                        total_benefit += base * (total_pct / 100.0)

                    # Distribuir beneficios por módulo
                    for module in module_selected:
                        benefit_module = 0.0
                        for aspect in module_profiles.get(module, []):
                            pct_module_aspect = module_benefits.get(module, {}).get(aspect, 0.0)
                            if aspect_totals.get(aspect, 0) > 0:
                                # Relación proporcional del beneficio total de ese aspecto atribuido a este módulo
                                a_lower = aspect.lower()
                                if any(k in a_lower for k in ["inv", "stock"]): base = inventory * (carrying_cost / 100.0)
                                elif any(k in a_lower for k in ["lab", "sga", "pers", "tra"]): base = sga_year
                                elif any(k in a_lower for k in ["marg", "prof", "log"]): base = base_oi_year
                                else: base = rev_year
                                
                                aspect_benefit = base * (aspect_totals[aspect] / 100.0)
                                benefit_module += aspect_benefit * (pct_module_aspect / aspect_totals[aspect])

                        investment_module = total_investment / len(module_selected) if module_selected else 0.0
                        cashflow_module = benefit_module - investment_module

                        module_results[module]["benefit"].append(round(benefit_module, 2))
                        module_results[module]["investment"].append(round(investment_module, 2))
                        module_results[module]["cashflow"].append(round(cashflow_module, 2))
                        cumulative_mod = (module_results[module]["cumulative"][-1] if module_results[module]["cumulative"] else 0.0) + cashflow_module
                        module_results[module]["cumulative"].append(round(cumulative_mod, 2))

                    yearly_total_benefit.append(round(total_benefit, 2))
                    yearly_total_investment.append(round(total_investment, 2))
                    cashflow_total = total_benefit - total_investment
                    yearly_total_cashflow.append(round(cashflow_total, 2))
                    cumulative_total += cashflow_total
                    yearly_total_cumulative.append(round(cumulative_total, 2))

                npv_value = npv(discount_rate / 100.0, yearly_total_cashflow)
                irr_value = irr(yearly_total_cashflow)
                payback = payback_period(yearly_total_cumulative)

                return {
                    "cliente": client_name,
                    "retailer_type": retailer_type,
                    "base_net_revenue": net_revenue,
                    "years": years,
                    "total_benefit": yearly_total_benefit,
                    "total_investment": yearly_total_investment,
                    "total_cashflow": yearly_total_cashflow,
                    "total_cumulative": yearly_total_cumulative,
                    "module_results": module_results,
                    "npv": round(npv_value, 2),
                    "irr": round(irr_value * 100.0, 2) if irr_value is not None else None,
                    "payback": payback,
                    "selection": module_selected,
                    "aspect_pcts": module_benefits
                }

            results = calc_manual_roi()
            
            # Guardar en session_state y mostrar reporte inmediatamente
            st.session_state.manual_results = results
            show_executive_report(results)

        if "manual_results" in st.session_state:
            st.write("---")
            st.success("Análisis completado. El detalle completo está disponible en el Reporte Ejecutivo pop-up.")
            if st.button("🔍 Abrir Reporte Ejecutivo"):
                show_executive_report(st.session_state.manual_results)

